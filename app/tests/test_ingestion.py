"""Tests for src/ingestion/parse_slides.py."""

import fitz
import pytest
from pptx import Presentation
from pptx.util import Inches

from src.ingestion.parse_slides import parse_pdf, parse_pptx


@pytest.fixture
def sample_pdf(tmp_path):
    path = tmp_path / "lecture.pdf"
    doc = fitz.open()
    page = doc.new_page()
    page.insert_text((72, 72), "Introduction to Neural Networks")
    doc.save(str(path))
    doc.close()
    return path


@pytest.fixture
def sample_pptx(tmp_path):
    path = tmp_path / "lecture.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Lecture 1"
    slide.placeholders[1].text = "Gradient descent minimizes the loss function."
    prs.save(str(path))
    return path


def test_parse_pdf_returns_list(sample_pdf):
    chunks = parse_pdf(str(sample_pdf))
    assert isinstance(chunks, list)
    assert len(chunks) >= 1


def test_parse_pdf_chunk_keys(sample_pdf):
    chunks = parse_pdf(str(sample_pdf))
    for chunk in chunks:
        assert "page" in chunk
        assert "text" in chunk
        assert "source" in chunk


def test_parse_pdf_page_is_int(sample_pdf):
    chunks = parse_pdf(str(sample_pdf))
    assert all(isinstance(c["page"], int) for c in chunks)


def test_parse_pdf_source_is_filename(sample_pdf):
    chunks = parse_pdf(str(sample_pdf))
    assert all(c["source"] == "lecture.pdf" for c in chunks)


def test_parse_pdf_text_contains_content(sample_pdf):
    chunks = parse_pdf(str(sample_pdf))
    combined = " ".join(c["text"] for c in chunks)
    assert "Neural Networks" in combined


def test_parse_pptx_returns_list(sample_pptx):
    chunks = parse_pptx(str(sample_pptx))
    assert isinstance(chunks, list)
    assert len(chunks) >= 1


def test_parse_pptx_chunk_keys(sample_pptx):
    chunks = parse_pptx(str(sample_pptx))
    for chunk in chunks:
        assert "page" in chunk
        assert "text" in chunk
        assert "source" in chunk


def test_parse_pptx_page_is_int(sample_pptx):
    chunks = parse_pptx(str(sample_pptx))
    assert all(isinstance(c["page"], int) for c in chunks)


def test_parse_pptx_source_is_filename(sample_pptx):
    chunks = parse_pptx(str(sample_pptx))
    assert all(c["source"] == "lecture.pptx" for c in chunks)


def test_parse_pptx_text_contains_content(sample_pptx):
    chunks = parse_pptx(str(sample_pptx))
    combined = " ".join(c["text"] for c in chunks)
    assert "Gradient descent" in combined
