"""Slide and PDF parsing utilities for lecture content ingestion."""

import sys
from pathlib import Path


def parse_pdf(path: str) -> list[dict]:
    """Parse a PDF file into a list of page chunks.

    Args:
        path: Absolute or relative path to the PDF file.

    Returns:
        A list of dicts with keys:
            page (int): 1-based page number.
            text (str): Extracted text content of the page.
            source (str): Filename of the source document.
    """
    import fitz  # pymupdf

    source = Path(path).name
    chunks = []
    with fitz.open(path) as doc:
        for i, page in enumerate(doc, start=1):
            text = page.get_text().strip()
            if text:
                chunks.append({"page": i, "text": text, "source": source})
    return chunks


def parse_pptx(path: str) -> list[dict]:
    """Parse a PPTX file into a list of slide chunks.

    Args:
        path: Absolute or relative path to the PPTX file.

    Returns:
        A list of dicts with keys:
            page (int): 1-based slide number.
            text (str): Concatenated text from all shapes on the slide.
            source (str): Filename of the source document.
    """
    from pptx import Presentation

    source = Path(path).name
    chunks = []
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides, start=1):
        texts = [
            shape.text_frame.text.strip()
            for shape in slide.shapes
            if shape.has_text_frame
        ]
        text = "\n".join(t for t in texts if t)
        if text:
            chunks.append({"page": i, "text": text, "source": source})
    return chunks


if __name__ == "__main__":
    import json

    if len(sys.argv) < 2:
        print("Usage: python parse_slides.py <file.pdf|file.pptx> [out_dir]")
        sys.exit(1)

    path = sys.argv[1]
    suffix = Path(path).suffix.lower()

    if suffix == ".pdf":
        chunks = parse_pdf(path)
    elif suffix in {".pptx", ".ppt"}:
        chunks = parse_pptx(path)
    else:
        print(f"Unsupported file type: {suffix}")
        sys.exit(1)

    print(f"{len(chunks)} chunks extracted from {Path(path).name}")

    if len(sys.argv) >= 3:
        out_dir = Path(sys.argv[2])
        out_dir.mkdir(parents=True, exist_ok=True)
        stem = Path(path).stem
        out_path = out_dir / f"{stem}.jsonl"
        with out_path.open("w", encoding="utf-8") as f:
            for chunk in chunks:
                f.write(json.dumps(chunk) + "\n")
        print(f"Saved → {out_path}")
